"""
PowerPoint generator using python-pptx.

PURPOSE:
This module generates the final PowerPoint report by filling in a template.
It takes the data extracted by excel_reader and replaces placeholders in the template.

KEY CONCEPTS:
1. PLACEHOLDERS: The template contains {{field_name}} tags that get replaced
   Example: {{month}} becomes "March", {{re_req}} becomes "45"

2. TWO TYPES OF REPLACEMENTS:
   a) Scalar replacements: {{month}}, {{re_req}}, etc.
   b) Table replacements: {{table:table_name}} creates dynamic tables

3. POWERPOINT QUIRK: PowerPoint often splits text across multiple "runs"
   (text fragments). We need special logic to handle placeholders that
   span multiple runs. Example: "{{mon" in run 1, "th}}" in run 2

KEY METHODS:
- generate_report(): Main entry point
- _replace_placeholders(): Finds and replaces all {{}} tags
- _populate_table(): Creates dynamic multi-row tables
- _replace_in_text_frame(): Handles the PowerPoint text run complexity
"""

import re
from copy import deepcopy  # For duplicating table rows
from datetime import datetime  # For timestamp in output filename
from pathlib import Path

from pptx import Presentation  # python-pptx library


class PowerPointWriter:
    """
    Generates PowerPoint reports by replacing placeholders in a template.
    
    WORKFLOW:
    1. Load a PowerPoint template file
    2. Find all {{placeholder}} tags in the template
    3. Replace them with data values from the Excel extraction
    4. Save the filled-in presentation
    
    This class handles both simple text replacements and complex table generation.
    """

    def __init__(self, config):
        """
        Initialize the PowerPoint writer.

        Args:
            config: Configuration dictionary containing:
                   - template_path: Path to the PowerPoint template
                   - output_dir: Where to save generated reports
        """
        self.config = config
        self.template_path = Path(config["template_path"])
        self.output_dir = Path(config.get("output_dir", "output"))

        # Ensure output directory exists
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_report(self, data):
        """
        Generate a PowerPoint report by replacing placeholders with data.
        
        This is the main method that orchestrates the report generation.
        It loads the template, replaces all placeholders, and saves the result.

        Args:
            data: Dictionary mapping field names to values
                 Scalar example: {"month": "March", "re_req": 45}
                 Table example: {"re_sup_cat": [{"re_cat": "Technical", "re_cat_count": 15}]}

        Returns:
            Path to the generated PowerPoint file
            
        Example:
            writer = PowerPointWriter(config)
            output = writer.generate_report({"month": "March", "re_req": 45})
            print(output)  # "output/demo_report_20240315_143022.pptx"
        """
        # STEP 1: Load the PowerPoint template
        prs = Presentation(self.template_path)

        # STEP 2: Replace all {{placeholder}} tags with actual data
        replacements_made = self._replace_placeholders(prs, data)

        print(f"  Made {replacements_made} replacements")

        # STEP 3: Generate output filename with timestamp
        # Format: YYYYMMDD_HHMMSS (e.g., 20240315_143022)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Include "demo" in filename if using demo config (helps identify test reports)
        config_path = self.config.get("config_path", "")
        if "demo" in config_path.lower():
            output_filename = f"demo_report_{timestamp}.pptx"
        else:
            output_filename = f"report_{timestamp}.pptx"

        output_path = self.output_dir / output_filename

        # STEP 4: Save the presentation
        prs.save(str(output_path))

        return output_path

    def _replace_placeholders(self, presentation, data):
        """
        Replace all {{placeholder}} tokens in the presentation with data values.
        
        This method searches through every slide, shape, and table in the presentation
        looking for {{field_name}} tags and replacing them with actual values.

        Args:
            presentation: python-pptx Presentation object (the loaded template)
            data: Dictionary of field names to values
                 Can include both scalar values and lists (for tables)

        Returns:
            Number of replacements made (for logging/debugging)
        """
        replacements = 0
        replaced_fields = set()  # Track which fields we've replaced

        # STEP 1: Separate table data from scalar data
        # Table data is identified by being a list instead of a single value
        table_data = {k: v for k, v in data.items() if isinstance(v, list)}
        scalar_data = {k: v for k, v in data.items() if not isinstance(v, list)}

        # STEP 2: Iterate through all slides and shapes
        for _slide_idx, slide in enumerate(presentation.slides):
            # Iterate through all shapes in the slide (text boxes, tables, etc.)
            for shape in slide.shapes:
                # Handle text shapes (text boxes, titles, etc.)
                if hasattr(shape, "text_frame"):
                    count, fields = self._replace_in_text_frame(
                        shape.text_frame, scalar_data
                    )
                    replacements += count
                    replaced_fields.update(fields)

                # Handle table shapes
                if shape.has_table:
                    # First, check if this is a DYNAMIC table (has {{table:name}} marker)
                    table_replacement = self._populate_table(shape.table, table_data)
                    if table_replacement:
                        # It was a dynamic table - we populated it
                        replacements += table_replacement
                    else:
                        # It's a static table - just replace placeholders in cells
                        for row in shape.table.rows:
                            for cell in row.cells:
                                count, fields = self._replace_in_text_frame(
                                    cell.text_frame, scalar_data
                                )
                                replacements += count
                                replaced_fields.update(fields)

        # STEP 3: Logging - show what was replaced
        print(f"  Replaced fields: {', '.join(sorted(replaced_fields))}")

        # STEP 4: Warning - show fields that WEREN'T replaced (might be typos)
        unreplaced = set(scalar_data.keys()) - replaced_fields
        if unreplaced:
            print(f"  ⚠ Unreplaced fields: {', '.join(sorted(unreplaced))}")

        return replacements

    def _populate_table(self, table, table_data):
        """
        Populate a table with dynamic rows if it contains a {{table:name}} placeholder.
        
        CONCEPT: Dynamic Tables
        Instead of manually creating each row in PowerPoint, we use a template approach:
        1. The template has a HEADER row (row 0) with column titles
        2. Row 1 is a TEMPLATE row with placeholders like {{re_cat}}, {{re_cat_count}}
        3. This method duplicates the template row for each data row
        4. It replaces the placeholders in each duplicated row with actual data
        
        Example Template Table:
        Row 0 (Header):  | Category | Count |
        Row 1 (Template): | {{re_cat}} | {{re_cat_count}} |
        
        Becomes:
        Row 0 (Header):  | Category | Count |
        Row 1 (Data):    | Technical | 15 |
        Row 2 (Data):    | Billing | 8 |
        Row 3 (Data):    | Access | 3 |

        Args:
            table: python-pptx Table object
            table_data: Dictionary of table names to lists of row data
                       Example: {"re_sup_cat": [{"re_cat": "Technical", "re_cat_count": 15}, ...]}

        Returns:
            Number of replacements made (int), or None if not a dynamic table
        """
        # VALIDATION: Check if table has at least 2 rows (header + template)
        if len(table.rows) < 2:
            return None

        # STEP 1: Detect if this is a dynamic table
        # Check if the first cell of row 1 (template row) contains {{table:table_name}}
        template_cell = table.rows[1].cells[0]
        template_cell_text = ""

        # Reconstruct full text from runs (PowerPoint splits text into fragments)
        for paragraph in template_cell.text_frame.paragraphs:
            for run in paragraph.runs:
                template_cell_text += run.text

        # Look for {{table:name}} pattern
        table_match = re.search(r"\{\{table:(\w+)\}\}", template_cell_text)

        # If no {{table:name}} found, this is not a dynamic table
        if not table_match:
            return None

        # Extract the table name from {{table:name}}
        table_name = table_match.group(1)  # e.g., "re_sup_cat"

        # VALIDATION: Check if we have data for this table
        if table_name not in table_data:
            print(f"  Warning: Table '{table_name}' not found in data")
            return 0

        rows_data = table_data[table_name]  # List of dicts, one dict per row

        if not rows_data:
            print(f"  Warning: Table '{table_name}' has no data rows")
            return 0

        # STEP 2: Extract template row information
        # Row 1 is the template that defines the structure
        template_row_idx = 1
        template_row = table.rows[template_row_idx]

        # Extract template cell properties and placeholders
        # We need to preserve the template structure for each data row
        template_cells = []
        for cell in template_row.cells:
            # Reconstruct full text from runs
            cell_text = ""
            for paragraph in cell.text_frame.paragraphs:
                cell_text += "".join(run.text for run in paragraph.runs)

            # Store both the text and the text_frame for reference
            cell_info = {"text": cell_text, "text_frame": cell.text_frame}
            template_cells.append(cell_info)

        # STEP 3: Determine how many rows we need
        # We need 1 header row + N data rows
        num_data_rows = len(rows_data)
        num_total_rows_needed = 1 + num_data_rows

        # STEP 4: Add more rows if needed by duplicating the template row
        # This uses XML manipulation to duplicate rows with formatting intact
        while len(table.rows) < num_total_rows_needed:
            # Access the underlying XML table element
            tbl = table._tbl
            # Deep copy the template row (preserves formatting, styles, etc.)
            tr = deepcopy(tbl.tr_lst[template_row_idx])
            # Append the duplicated row to the table
            tbl.append(tr)

        # STEP 5: Populate each data row
        # Start from row 1 (preserve row 0 as header)
        replacements = 0
        for data_idx, row_data in enumerate(rows_data):
            row_idx = data_idx + 1  # Row 1, 2, 3, etc. (Skip row 0 which is header)
            row = table.rows[row_idx]

            # For each cell in this row
            for col_idx, cell in enumerate(row.cells):
                # Get the template text for this column
                # Example: "{{re_cat}}" or "{{re_cat_count}}"
                template_text = template_cells[col_idx]["text"]

                # Replace placeholders in the template text with actual values
                new_text = template_text
                for field_name, value in row_data.items():
                    placeholder = f"{{{{{field_name}}}}}"  # e.g., "{{re_cat}}"
                    if placeholder in new_text:
                        new_text = new_text.replace(placeholder, str(value))
                        replacements += 1

                # Also remove the {{table:name}} marker from the first cell
                new_text = re.sub(r"\{\{table:\w+\}\}", "", new_text).strip()

                # STEP 6: Set the cell text while preserving formatting
                # Note: cell.text = new_text would lose all formatting (font, color, etc.)
                # Instead, we modify runs to preserve formatting
                try:
                    if cell.text_frame.paragraphs:
                        paragraph = cell.text_frame.paragraphs[0]
                        if paragraph.runs:
                            # Set text in first run (preserves its formatting)
                            first_run = paragraph.runs[0]
                            first_run.text = new_text
                            # Remove all other runs safely
                            while len(paragraph.runs) > 1:
                                try:
                                    # Access the underlying XML element to remove runs
                                    p = paragraph._element
                                    p.remove(paragraph.runs[-1]._element)
                                except (AttributeError, IndexError):
                                    # If _element access fails, break and use fallback
                                    break
                        else:
                            # No runs, fall back to simple assignment
                            cell.text = new_text
                    else:
                        # No paragraphs, fall back to simple assignment
                        cell.text = new_text
                except Exception:
                    # If anything fails, use simple assignment (may lose formatting)
                    cell.text = new_text

        print(
            f"  Populated table '{table_name}' with {num_data_rows} rows ({replacements} replacements)"
        )

        return replacements

    def _replace_in_text_frame(self, text_frame, data):
        """
        Replace placeholders in a text frame (text box or cell).
        
        THE POWERPOINT TEXT RUN PROBLEM:
        PowerPoint doesn't store text as a single string. Instead, it stores text
        as multiple "runs" (text fragments), where each run can have different formatting.
        
        PROBLEM: PowerPoint may SPLIT a placeholder across multiple runs!
        
        Example - what you see in PowerPoint:
            "The current month is {{month}}"
        
        What PowerPoint actually stores internally:
            Run 1: "The current month is {{mon"
            Run 2: "th}}"
        
        WHY THIS HAPPENS:
        - User edited the text multiple times
        - Copy/paste from different sources
        - PowerPoint's internal text handling quirks
        
        OUR SOLUTION:
        1. Reassemble all runs into the full original text
        2. Do all replacements on the full text
        3. Put the result in the first run (preserving its formatting)
        4. Clear the other runs
        
        This ensures we find ALL placeholders and preserve formatting.

        Args:
            text_frame: python-pptx TextFrame object (from a text box or cell)
            data: Dictionary of field names to values {"month": "March", ...}

        Returns:
            Tuple of (number of replacements made, set of field names that were replaced)
            Example: (2, {"month", "re_req"})
        """
        replacements = 0
        replaced_fields = set()

        # Process each paragraph in the text frame
        # (Most text frames have just one paragraph, but PowerPoint allows multiple)
        for paragraph in text_frame.paragraphs:
            # Skip empty paragraphs (no runs = no text)
            if not paragraph.runs:
                continue

            # STEP 1: Reconstruct the full paragraph text by joining all runs
            # This handles the case where a placeholder is split across runs
            full_text = "".join(run.text for run in paragraph.runs)

            # STEP 2: Check if there are ANY placeholders in this paragraph
            # This is an optimization - don't process paragraphs with no placeholders
            has_placeholder = False
            for field_name in data:
                placeholder = f"{{{{{field_name}}}}}"  # e.g., "{{month}}"
                if placeholder in full_text:
                    has_placeholder = True
                    break

            # If no placeholders found, skip to next paragraph
            if not has_placeholder:
                continue

            # STEP 3: Replace all placeholders in the full text
            # Now we can find placeholders that were split across runs
            new_text = full_text
            for field_name, value in data.items():
                placeholder = f"{{{{{field_name}}}}}"
                if placeholder in new_text:
                    # Replace the placeholder with the actual value
                    new_text = new_text.replace(placeholder, str(value))
                    replacements += 1
                    replaced_fields.add(field_name)

            # STEP 4: Put all text in the first run and clear the rest
            # WHY: This preserves the first run's formatting (font, size, color, etc.)
            # The first run usually has the desired formatting
            first_run = paragraph.runs[0]
            first_run.text = new_text

            # Clear text from all other runs
            # NOTE: We clear instead of removing because removing is complex in PowerPoint XML
            # Clearing is simpler and achieves the same result
            for run in paragraph.runs[1:]:
                run.text = ""

        return replacements, replaced_fields
