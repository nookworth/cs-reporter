# CS Reporter
# PowerPoint writer with chart support (V3).
# Extends V2 functionality with automatic chart generation.
# **Date:** 2026-02-20
# **Status:** ✅ Completed
# **Component:** src

import logging

from pptx import Presentation

from src.chart_utils_v3 import (
    add_category_pie_chart,
    add_comparison_chart,
    add_top_orgs_chart,
)
from src.ppt_writer import PowerPointWriter


class PowerPointWriterV3(PowerPointWriter):
    """
    Enhanced PowerPoint writer with chart generation.
    Inherits from base PowerPointWriter and adds chart capabilities.
    """

    def generate_report(self, data: dict) -> str:
        """Generate a PowerPoint report with charts from the data."""
        output_path = super().generate_report(data)
        self._add_charts_to_report(output_path, data)

        return output_path

    def _add_charts_to_report(self, ppt_path: str, data: dict) -> None:
        """Add visual analytics charts to the PowerPoint report."""
        try:
            prs = Presentation(ppt_path)

            # Add a new slide for charts
            try:
                blank_layout = prs.slide_layouts[-1]  # Use last layout
            except Exception:
                blank_layout = prs.slide_layouts[0]  # Fallback to first layout

            chart_slide = prs.slides.add_slide(blank_layout)

            # Add title
            title_box = chart_slide.shapes.add_textbox(
                left=0, top=0, width=prs.slide_width, height=int(prs.slide_height * 0.1)
            )
            title_frame = title_box.text_frame
            title_frame.text = "📊 Visual Analytics"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = int(prs.slide_height * 0.05)
            title_para.font.bold = True

            # Add comparison chart (top left)
            try:
                add_comparison_chart(
                    chart_slide, data, left=0.5, top=1.5, width=4.5, height=3
                )
            except Exception as e:
                logging.warning(f"Could not add comparison chart: {e}")

            # Add category pie chart (top right)
            try:
                retail_cats = data.get("re_sup_cat", [])
                if retail_cats:
                    add_category_pie_chart(
                        chart_slide,
                        retail_cats,
                        "Retail Categories",
                        left=5.5,
                        top=1.5,
                        width=4,
                        height=3,
                    )
            except Exception as e:
                logging.warning(f"Could not add category chart: {e}")

            # Add top orgs chart (bottom)
            try:
                top_orgs = data.get("top_orgs", [])
                if top_orgs:
                    add_top_orgs_chart(
                        chart_slide, top_orgs, left=1.5, top=5, width=7, height=2.5
                    )
            except Exception as e:
                logging.warning(f"Could not add top orgs chart: {e}")

            # Save the presentation
            prs.save(ppt_path)
            print("  [OK] Added charts to presentation")

        except Exception as e:
            logging.warning(f"Could not add charts: {e}")
            # Don't fail the whole report if charts fail
