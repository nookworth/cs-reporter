"""
PowerPoint generator using python-pptx.
Replaces placeholders in a template with data from Excel.
"""

from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
import re


class PowerPointWriter:
    """Generates PowerPoint reports by replacing placeholders in a template."""

    def __init__(self, config):
        """
        Initialize the PowerPoint writer.

        Args:
            config: Configuration dictionary with template path
        """
        self.config = config
        self.template_path = Path(config['template_path'])
        self.output_dir = Path(config.get('output_dir', 'output'))

        # Ensure output directory exists
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_report(self, data):
        """
        Generate a PowerPoint report by replacing placeholders with data.

        Args:
            data: Dictionary mapping field names to values

        Returns:
            Path to the generated PowerPoint file
        """
        # Load the template
        prs = Presentation(self.template_path)

        # Replace placeholders in all slides
        replacements_made = self._replace_placeholders(prs, data)

        print(f"  Made {replacements_made} replacements")

        # Generate output filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"report_{timestamp}.pptx"
        output_path = self.output_dir / output_filename

        # Save the presentation
        prs.save(str(output_path))

        return output_path

    def _replace_placeholders(self, presentation, data):
        """
        Replace all {{placeholder}} tokens in the presentation with data values.

        Args:
            presentation: python-pptx Presentation object
            data: Dictionary of field names to values (includes tables as lists)

        Returns:
            Number of replacements made
        """
        replacements = 0
        replaced_fields = set()

        # Separate table data from scalar data
        table_data = {k: v for k, v in data.items() if isinstance(v, list)}
        scalar_data = {k: v for k, v in data.items() if not isinstance(v, list)}

        # Iterate through all slides
        for slide_idx, slide in enumerate(presentation.slides):
            # Iterate through all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    count, fields = self._replace_in_text_frame(
                        shape.text_frame, scalar_data
                    )
                    replacements += count
                    replaced_fields.update(fields)

                # Check tables
                if shape.has_table:
                    # Check if this table has a dynamic table placeholder
                    table_replacement = self._populate_table(shape.table, table_data)
                    if table_replacement:
                        replacements += table_replacement
                    else:
                        # Regular cell replacement for static tables
                        for row in shape.table.rows:
                            for cell in row.cells:
                                count, fields = self._replace_in_text_frame(
                                    cell.text_frame, scalar_data
                                )
                                replacements += count
                                replaced_fields.update(fields)

        # Show which fields were replaced
        print(f"  Replaced fields: {', '.join(sorted(replaced_fields))}")

        # Show which fields were NOT replaced (might still be placeholders)
        unreplaced = set(scalar_data.keys()) - replaced_fields
        if unreplaced:
            print(f"  ⚠ Unreplaced fields: {', '.join(sorted(unreplaced))}")

        return replacements

    def _populate_table(self, table, table_data):
        """
        Populate a table with dynamic rows if it contains a {{table:name}} placeholder.

        Table structure:
        - Row 0: Header row (preserved as-is)
        - Row 1: Template row with placeholders (used to create data rows)

        Args:
            table: python-pptx Table object
            table_data: Dictionary of table names to lists of row data

        Returns:
            Number of replacements made, or None if not a dynamic table
        """
        # Check if table has at least 2 rows (header + template)
        if len(table.rows) < 2:
            return None

        # Check if the first cell of the SECOND row contains {{table:table_name}}
        # (First row is the header)
        template_cell = table.rows[1].cells[0]
        template_cell_text = ''

        for paragraph in template_cell.text_frame.paragraphs:
            for run in paragraph.runs:
                template_cell_text += run.text

        table_match = re.search(r'\{\{table:(\w+)\}\}', template_cell_text)

        if not table_match:
            return None

        table_name = table_match.group(1)

        if table_name not in table_data:
            print(f"  Warning: Table '{table_name}' not found in data")
            return 0

        rows_data = table_data[table_name]

        if not rows_data:
            print(f"  Warning: Table '{table_name}' has no data rows")
            return 0

        # The second row (index 1) is the template
        template_row_idx = 1
        template_row = table.rows[template_row_idx]

        # Extract template cell properties and placeholders
        template_cells = []
        for cell in template_row.cells:
            # Reconstruct full text from runs (PowerPoint splits text)
            cell_text = ''
            for paragraph in cell.text_frame.paragraphs:
                cell_text += ''.join(run.text for run in paragraph.runs)

            cell_info = {
                'text': cell_text,
                'text_frame': cell.text_frame
            }
            template_cells.append(cell_info)

        # Determine how many total rows we need (header + data rows)
        num_data_rows = len(rows_data)
        num_total_rows_needed = 1 + num_data_rows  # 1 header + N data rows

        # Add more rows if needed by duplicating the template row
        while len(table.rows) < num_total_rows_needed:
            # Duplicate the template row using the underlying XML element
            tbl = table._tbl
            tr = deepcopy(tbl.tr_lst[template_row_idx])
            tbl.append(tr)

        # Populate each data row (starting from row 1, preserving row 0 as header)
        replacements = 0
        for data_idx, row_data in enumerate(rows_data):
            row_idx = data_idx + 1  # Skip row 0 (header)
            row = table.rows[row_idx]

            for col_idx, cell in enumerate(row.cells):
                # Get template text for this column
                template_text = template_cells[col_idx]['text']

                # Replace placeholders in the template text
                new_text = template_text
                for field_name, value in row_data.items():
                    placeholder = f"{{{{{field_name}}}}}"
                    if placeholder in new_text:
                        new_text = new_text.replace(placeholder, str(value))
                        replacements += 1

                # Also remove the {{table:name}} marker from first cell
                new_text = re.sub(r'\{\{table:\w+\}\}', '', new_text).strip()

                # Set the cell text while preserving formatting
                # Instead of cell.text = new_text (which clears formatting),
                # we modify the first run and remove others
                if cell.text_frame.paragraphs:
                    paragraph = cell.text_frame.paragraphs[0]
                    if paragraph.runs:
                        # Set text in first run (preserves its formatting)
                        first_run = paragraph.runs[0]
                        first_run.text = new_text
                        # Remove all other runs
                        for _ in range(len(paragraph.runs) - 1):
                            p = paragraph._element
                            p.remove(paragraph.runs[-1]._element)
                    else:
                        # No runs, fall back to simple assignment
                        cell.text = new_text
                else:
                    # No paragraphs, fall back to simple assignment
                    cell.text = new_text

        print(f"  Populated table '{table_name}' with {num_data_rows} rows ({replacements} replacements)")

        return replacements

    def _replace_in_text_frame(self, text_frame, data):
        """
        Replace placeholders in a text frame.

        Handles PowerPoint's tendency to split placeholders across multiple runs.

        Args:
            text_frame: python-pptx TextFrame object
            data: Dictionary of field names to values

        Returns:
            Tuple of (number of replacements made, set of field names that were replaced)
        """
        replacements = 0
        replaced_fields = set()

        for paragraph in text_frame.paragraphs:
            # Skip empty paragraphs
            if not paragraph.runs:
                continue

            # Get the full paragraph text
            full_text = ''.join(run.text for run in paragraph.runs)

            # Check if there are any placeholders in this paragraph
            has_placeholder = False
            for field_name in data.keys():
                placeholder = f"{{{{{field_name}}}}}"
                if placeholder in full_text:
                    has_placeholder = True
                    break

            if not has_placeholder:
                continue

            # Replace all placeholders in the full text
            new_text = full_text
            for field_name, value in data.items():
                placeholder = f"{{{{{field_name}}}}}"
                if placeholder in new_text:
                    new_text = new_text.replace(placeholder, str(value))
                    replacements += 1
                    replaced_fields.add(field_name)

            # Clear all runs and set the new text in a single run
            # This preserves the first run's formatting
            first_run = paragraph.runs[0]
            first_run.text = new_text

            # Remove all other runs
            for _ in range(len(paragraph.runs) - 1):
                p = paragraph._element
                p.remove(paragraph.runs[-1]._element)

        return replacements, replaced_fields
