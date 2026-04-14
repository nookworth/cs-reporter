# CS Reporter
# Chart utilities for adding charts to PowerPoint presentations.
# **Date:** 2026-02-20
# **Status:** ✅ Completed
# **Component:** src

from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Inches


def add_comparison_chart(
    slide: Slide, data: dict, left: float, top: float, width: float, height: float
) -> Chart:
    """Add a comparison bar chart comparing current and previous month metrics."""
    chart_data = CategoryChartData()
    chart_data.categories = ["Tickets", "Good Ratings", "Resolution Time"]

    # Current month data
    current_values = (
        data.get("re_req", 0) + data.get("su_req", 0),
        data.get("re_sat", 0) + data.get("su_sat", 0),
        data.get("re_reso", 0),
    )

    # Previous month data
    previous_values = (
        data.get("re_prev_req", 0) + data.get("su_prev_req", 0),
        data.get("re_prev_sat", 0) + data.get("su_prev_sat", 0),
        data.get("re_prev_reso", 0),
    )

    chart_data.add_series("Current Month", current_values)
    chart_data.add_series("Previous Month", previous_values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = 2  # Right

    return chart


def add_category_pie_chart(
    slide: Slide,
    table_data: list[dict],
    title: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> Chart | None:
    """Add a pie chart showing top categories from table data."""
    if not table_data:
        return None

    chart_data = CategoryChartData()

    # Extract categories and counts
    categories = []
    values = []
    for item in table_data[:5]:  # Top 5 categories
        # Handle different field name patterns
        cat_name = (
            item.get("re_cat") or item.get("su_cat") or item.get("value", "Unknown")
        )
        cat_count = (
            item.get("re_cat_count") or item.get("su_cat_count") or item.get("count", 0)
        )
        categories.append(cat_name)
        values.append(cat_count)

    chart_data.categories = categories
    chart_data.add_series(title, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = 2  # Right

    return chart


def add_top_orgs_chart(
    slide: Slide,
    table_data: list[dict],
    left: float,
    top: float,
    width: float,
    height: float,
) -> Chart | None:
    """Add a bar chart showing top organizations by ticket count."""
    if not table_data:
        return None

    chart_data = CategoryChartData()

    # Extract org names and counts
    org_names = []
    org_counts = []
    for item in table_data[:5]:  # Top 5 orgs
        org_names.append(item.get("org_name", "Unknown"))
        org_counts.append(item.get("org_count", 0))

    chart_data.categories = org_names
    chart_data.add_series("Tickets", org_counts)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart

    chart.has_legend = False

    return chart


def calculate_percentage_change(current: float, previous: float) -> float:
    """Calculate percentage change between current and previous values."""
    if previous == 0:
        return 0
    return round(((current - previous) / previous) * 100, 1)


def get_chart_summary(data: dict) -> dict[str, float]:
    """Generate summary metrics for charts from raw data."""
    total_current = data.get("re_req", 0) + data.get("su_req", 0)
    total_previous = data.get("re_prev_req", 0) + data.get("su_prev_req", 0)

    return {
        "total_tickets_current": total_current,
        "total_tickets_previous": total_previous,
        "ticket_change_pct": calculate_percentage_change(total_current, total_previous),
        "avg_resolution_current": data.get("re_reso", 0),
        "avg_resolution_previous": data.get("re_prev_reso", 0),
        "resolution_change_pct": calculate_percentage_change(
            data.get("re_reso", 0), data.get("re_prev_reso", 0)
        ),
    }
